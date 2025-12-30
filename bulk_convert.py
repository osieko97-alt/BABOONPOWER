"""
Bulk Document to HTML Converter

This script scans the workspace for PDF, DOCX, and PPTX files, extracts their content, and generates HTML files for each document. It also creates an index.html linking to all generated pages.

Dependencies:
- python-docx
- PyPDF2
- python-pptx

Usage:
1. Install dependencies:
   pip install python-docx PyPDF2 python-pptx
2. Run the script:
   python bulk_convert.py
"""

import os
import glob
from docx import Document
from pptx import Presentation
from PyPDF2 import PdfReader

WORKSPACE = os.path.dirname(os.path.abspath(__file__))
GENERATED_DIR = os.path.join(WORKSPACE, "generated")
ASSETS_DIR = os.path.join(GENERATED_DIR, "assets")

os.makedirs(GENERATED_DIR, exist_ok=True)
os.makedirs(ASSETS_DIR, exist_ok=True)

def extract_docx(path):
    doc = Document(path)
    text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
    title = os.path.basename(path)
    return title, text

def extract_pdf(path):
    reader = PdfReader(path)
    text = "\n".join([page.extract_text() or "" for page in reader.pages])
    title = os.path.basename(path)
    return title, text

def extract_pptx(path):
    prs = Presentation(path)
    slides = []
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        slides.append("<br>".join(slide_text))
    text = "<hr>".join(slides)
    title = os.path.basename(path)
    return title, text

def save_html(title, content, filename):
    html = f"""
    <html>
    <head><title>{title}</title></head>
    <body>
    <h1>{title}</h1>
    <div>{content.replace('\n', '<br>')}</div>
    </body>
    </html>
    """
    with open(filename, "w", encoding="utf-8") as f:
        f.write(html)

def main():
    files = []
    files += glob.glob(os.path.join(WORKSPACE, "**/*.pdf"), recursive=True)
    files += glob.glob(os.path.join(WORKSPACE, "**/*.docx"), recursive=True)
    files += glob.glob(os.path.join(WORKSPACE, "**/*.pptx"), recursive=True)
    index_entries = []
    for path in files:
        ext = os.path.splitext(path)[1].lower()
        if ext == ".docx":
            title, content = extract_docx(path)
        elif ext == ".pdf":
            title, content = extract_pdf(path)
        elif ext == ".pptx":
            title, content = extract_pptx(path)
        else:
            continue
        html_name = os.path.splitext(os.path.basename(path))[0] + ".html"
        html_path = os.path.join(GENERATED_DIR, html_name)
        save_html(title, content, html_path)
        index_entries.append(f'<li><a href="{html_name}">{title}</a></li>')
        # Copy original to assets
        asset_path = os.path.join(ASSETS_DIR, os.path.basename(path))
        if not os.path.exists(asset_path):
            try:
                import shutil
                shutil.copy2(path, asset_path)
            except Exception:
                pass
    # Create index.html
    index_html = f"""
    <html>
    <head><title>Document Index</title></head>
    <body>
    <h1>All Documents</h1>
    <ul>
    {''.join(index_entries)}
    </ul>
    </body>
    </html>
    """
    with open(os.path.join(GENERATED_DIR, "index.html"), "w", encoding="utf-8") as f:
        f.write(index_html)

if __name__ == "__main__":
    main()
