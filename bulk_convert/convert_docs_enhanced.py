#!/usr/bin/env python3
"""
Enhanced document converter with OCR + .doc support.
- Extracts text from PDF, DOCX, PPTX, DOC files.
- Uses OCR (pytesseract) for scanned PDFs.
- Converts .doc to .docx via LibreOffice headless.
- Generates one HTML per document.
Usage:
  python convert_docs_enhanced.py --src <source-folder> --out <output-folder>
"""
import argparse
import os
import shutil
import sys
import html
import subprocess
import tempfile
from pathlib import Path
from datetime import datetime

try:
    from docx import Document as DocxDocument
except Exception:
    DocxDocument = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    from PyPDF2 import PdfReader
except Exception:
    PdfReader = None

try:
    import pytesseract
    from PIL import Image
    PYTESSERACT_AVAILABLE = True
except Exception:
    PYTESSERACT_AVAILABLE = False

try:
    from pdf2image import convert_from_path
    PDF2IMAGE_AVAILABLE = True
except Exception:
    PDF2IMAGE_AVAILABLE = False


SUPPORTED = ['.pdf', '.docx', '.pptx', '.doc']

CSS = '''
body{font-family:Segoe UI,Arial,Helvetica,sans-serif;margin:24px}
header{border-bottom:1px solid #ddd;margin-bottom:12px;padding-bottom:8px}
.title{font-size:1.6rem;font-weight:700}
.meta{color:#666;font-size:0.9rem}
.content{white-space:pre-wrap;margin-top:16px;line-height:1.5}
.item{margin:8px 0;padding:10px;border:1px solid #f0f0f0;border-radius:6px}
.warning{color:#d9534f;font-weight:600}
'''


def safe_filename(s: str) -> str:
    keep = "-_ .()[]" + "0123456789abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ"
    s2 = ''.join(c if c in keep else '_' for c in s)
    s2 = s2.strip(' _.-')
    if not s2:
        s2 = 'doc'
    return s2


def extract_text_docx(path: Path) -> (str, str):
    if DocxDocument is None:
        return '', ''
    try:
        doc = DocxDocument(path)
        props = doc.core_properties
        title = props.title or path.stem
        parts = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                parts.append(text)
        return title, '\n\n'.join(parts)
    except Exception as e:
        return path.stem, f"[Error extracting DOCX: {e}]"


def extract_text_pptx(path: Path) -> (str, str):
    if Presentation is None:
        return '', ''
    try:
        prs = Presentation(str(path))
        props = prs.core_properties
        title = props.title or path.stem
        slides = []
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    t = shape.text.strip()
                    if t:
                        slide_text.append(t)
            if slide_text:
                slides.append('\n'.join(slide_text))
        content = '\n\n'.join(f"Slide {i+1}:\n{slides[i]}" for i in range(len(slides)))
        return title, content
    except Exception as e:
        return path.stem, f"[Error extracting PPTX: {e}]"


def ocr_pdf(path: Path) -> str:
    """Use pytesseract + pdf2image to OCR a PDF."""
    if not PYTESSERACT_AVAILABLE or not PDF2IMAGE_AVAILABLE:
        return None
    try:
        # Convert PDF to images (limit to first 10 pages for speed)
        images = convert_from_path(str(path), first_page=1, last_page=10, dpi=200)
        if not images:
            return None
        ocr_text = []
        for i, img in enumerate(images):
            try:
                text = pytesseract.image_to_string(img)
                if text.strip():
                    ocr_text.append(f"[Page {i+1}]\n{text}")
            except Exception:
                pass
        return '\n\n'.join(ocr_text) if ocr_text else None
    except Exception:
        return None


def extract_text_pdf(path: Path) -> (str, str):
    if PdfReader is None:
        return '', ''
    try:
        reader = PdfReader(str(path))
        info = reader.metadata
        title = None
        if info and hasattr(info, 'title'):
            title = info.title
        title = title or path.stem
        pages = []
        for p in reader.pages:
            try:
                t = p.extract_text()
            except Exception:
                t = None
            if t and t.strip():
                pages.append(t)
        content = '\n\n'.join(pages) if pages else None
        # If no text extracted, try OCR (scanned PDF)
        if not content or len((content or '').strip()) < 50:
            ocr_content = ocr_pdf(path)
            if ocr_content:
                content = ocr_content
            else:
                content = "[PDF: No text extracted. (Scanned PDF without OCR available.)]"
        return title, content
    except Exception as e:
        return path.stem, f"[Error extracting PDF: {e}]"


def convert_doc_to_docx(path: Path, tmpdir: Path) -> Path:
    """Convert .doc to .docx using LibreOffice headless."""
    try:
        subprocess.run([
            'soffice', '--headless', '--convert-to', 'docx',
            '--outdir', str(tmpdir),
            str(path)
        ], check=True, capture_output=True, timeout=30)
        docx_path = tmpdir / (path.stem + '.docx')
        if docx_path.exists():
            return docx_path
    except Exception:
        pass
    return None


def extract_text_doc(path: Path) -> (str, str):
    """Extract from .doc by converting to .docx first."""
    if DocxDocument is None:
        return '', ''
    with tempfile.TemporaryDirectory() as tmpdir:
        tmpdir = Path(tmpdir)
        docx_path = convert_doc_to_docx(path, tmpdir)
        if docx_path:
            return extract_text_docx(docx_path)
    return path.stem, "[.doc file: LibreOffice not available for conversion.]"


def write_html(out_path: Path, title: str, meta: dict, content: str, original_rel: str):
    safe_title = html.escape(title)
    meta_lines = '\n'.join(f"<div><strong>{html.escape(k)}:</strong> {html.escape(str(v))}</div>" for k, v in meta.items() if v)
    body = f"""<!doctype html>
<html lang="en">
<head>
<meta charset="utf-8" />
<meta name="viewport" content="width=device-width,initial-scale=1" />
<title>{safe_title}</title>
<style>{CSS}</style>
</head>
<body>
<header>
<div class="title">{safe_title}</div>
<div class="meta">{meta_lines}</div>
<div class="meta">Original file: <a href="{html.escape(original_rel)}">{html.escape(os.path.basename(original_rel))}</a></div>
</header>
<main class="content"><pre>{html.escape(content)}</pre></main>
</body>
</html>
"""
    out_path.write_text(body, encoding='utf-8')


def scan_and_convert(src: Path, out: Path):
    src = src.resolve()
    out = out.resolve()
    assets = out / 'assets'
    pages = out / 'pages'
    out.mkdir(parents=True, exist_ok=True)
    assets.mkdir(parents=True, exist_ok=True)
    pages.mkdir(parents=True, exist_ok=True)

    index_items = []
    count = 0

    for root, dirs, files in os.walk(src):
        for f in files:
            p = Path(root) / f
            ext = p.suffix.lower()
            if ext not in SUPPORTED:
                continue
            rel_src = p.relative_to(src)
            print(f"[{count+1}] Processing {rel_src}")
            count += 1

            # copy original into assets
            target_asset = assets / safe_filename(str(rel_src))
            try:
                shutil.copy2(p, target_asset)
            except Exception:
                target_asset = assets / safe_filename(p.name)
                shutil.copy2(p, target_asset)

            if ext == '.docx':
                title, content = extract_text_docx(p)
            elif ext == '.pptx':
                title, content = extract_text_pptx(p)
            elif ext == '.pdf':
                title, content = extract_text_pdf(p)
            elif ext == '.doc':
                title, content = extract_text_doc(p)
            else:
                title, content = p.stem, '[Unsupported file type]'

            meta = {
                'source_path': str(rel_src),
                'converted_at': datetime.utcnow().isoformat() + 'Z',
                'size_bytes': p.stat().st_size,
                'file_type': ext.upper()
            }
            page_name = safe_filename(title) + '_' + str(abs(hash(str(rel_src))))[:8] + '.html'
            page_path = pages / page_name
            original_rel = os.path.relpath(target_asset, page_path.parent)
            write_html(page_path, title, meta, content or '[No text extracted]', original_rel)

            index_items.append({
                'title': title,
                'page': os.path.relpath(page_path, out),
                'original': os.path.relpath(target_asset, out),
                'snippet': (content or '')[:250].replace('\n', ' '),
                'file_type': ext.upper()
            })

    # write index
    idx_lines = ["<!doctype html>", "<html><head><meta charset=\"utf-8\"><title>Document Index</title>", f"<style>{CSS}</style>", "</head><body>", "<h1>Document Library</h1>", f"<p>Total documents converted: {len(index_items)}</p>", "<div>"]
    for it in index_items:
        idx_lines.append(f"<div class='item'><a href='{html.escape(it['page'])}'><strong>{html.escape(it['title'])}</strong></a> <span style='color:#999;font-size:0.9rem'>[{it['file_type']}]</span><div class='meta'>Original: <a href='{html.escape(it['original'])}'>{html.escape(os.path.basename(it['original']))}</a></div><div>{html.escape(it['snippet'])}...</div></div>")
    idx_lines.append("</div></body></html>")
    (out / 'index.html').write_text('\n'.join(idx_lines), encoding='utf-8')
    print(f"\n✓ Converted {len(index_items)} documents to {out}")
    print(f"  - Pages: {out / 'pages'}")
    print(f"  - Assets: {out / 'assets'}")
    print(f"  - Index: {out / 'index.html'}")


def main():
    parser = argparse.ArgumentParser(description='Convert documents (PDF, DOCX, PPTX, DOC) to HTML')
    parser.add_argument('--src', help='Source folder to scan', default='.')
    parser.add_argument('--out', help='Output folder', default='generated')
    args = parser.parse_args()

    src = Path(args.src)
    out = Path(args.out)
    if not src.exists():
        print(f'Source folder does not exist: {src}')
        sys.exit(1)

    scan_and_convert(src, out)


if __name__ == '__main__':
    main()
