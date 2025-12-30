#!/usr/bin/env python3
"""
Fast document converter (text-only, no OCR).
Converts PDF, DOCX, PPTX, DOC files to HTML pages.
"""
import argparse
import os
import shutil
import sys
import html
import subprocess
import tempfile
from pathlib import Path
from datetime import datetime

try:
    from docx import Document as DocxDocument
except:
    DocxDocument = None

try:
    from pptx import Presentation
except:
    Presentation = None

try:
    from PyPDF2 import PdfReader
except:
    PdfReader = None


SUPPORTED = ['.pdf', '.docx', '.pptx', '.doc']
CSS = 'body{font-family:Segoe UI,Arial;margin:20px}h1{font-size:1.5rem}pre{background:#f5f5f5;padding:12px;overflow:auto}.meta{color:#999;font-size:0.9rem}.index{display:grid;gap:12px}.item{border:1px solid #ddd;padding:10px;border-radius:4px}'


def safe_filename(s):
    keep = "-_ .()[]" + "0123456789abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ"
    s2 = ''.join(c if c in keep else '_' for c in s)
    return s2.strip(' _.') or 'doc'


def extract_docx(path):
    if not DocxDocument:
        return path.stem, ''
    try:
        doc = DocxDocument(path)
        title = doc.core_properties.title or path.stem
        text = '\n\n'.join(p.text for p in doc.paragraphs if p.text.strip())
        return title, text
    except:
        return path.stem, '[Error reading DOCX]'


def extract_pptx(path):
    if not Presentation:
        return path.stem, ''
    try:
        prs = Presentation(str(path))
        title = prs.core_properties.title or path.stem
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            slide_text = '\n'.join(s.text for s in slide.shapes if hasattr(s, 'text') and s.text)
            if slide_text:
                slides.append(f"--- Slide {i} ---\n{slide_text}")
        return title, '\n\n'.join(slides)
    except:
        return path.stem, '[Error reading PPTX]'


def extract_pdf(path):
    if not PdfReader:
        return path.stem, ''
    try:
        reader = PdfReader(str(path))
        title = (reader.metadata.title if reader.metadata else None) or path.stem
        text = ''
        for p in reader.pages:
            try:
                text += p.extract_text() + '\n\n'
            except:
                pass
        return title, text or '[No text found in PDF]'
    except:
        return path.stem, '[Error reading PDF]'


def convert_doc(path, tmp):
    try:
        subprocess.run(['soffice', '--headless', '--convert-to', 'docx', '--outdir', str(tmp), str(path)], 
                      capture_output=True, timeout=20)
        docx = tmp / (path.stem + '.docx')
        return docx if docx.exists() else None
    except:
        return None


def extract_doc(path):
    if not DocxDocument:
        return path.stem, '[LibreOffice not available]'
    with tempfile.TemporaryDirectory() as tmp:
        docx = convert_doc(path, Path(tmp))
        if docx:
            return extract_docx(docx)
    return path.stem, '[Could not convert .doc]'


def write_html(out_path, title, meta, content, asset_link):
    html_text = f'''<!DOCTYPE html>
<html><head><meta charset="utf-8"><meta name="viewport" content="width=device-width">
<title>{html.escape(title)}</title><style>{CSS}</style></head><body>
<h1>{html.escape(title)}</h1>
<div class="meta">{' | '.join(f"{k}: {html.escape(str(v))}" for k,v in meta.items())}</div>
<p><a href="{html.escape(asset_link)}">📄 Original File</a></p>
<pre>{html.escape(content[:10000])}</pre>
</body></html>'''
    out_path.write_text(html_text, encoding='utf-8')


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument('--src', default='.')
    parser.add_argument('--out', default='generated')
    args = parser.parse_args()

    src = Path(args.src).resolve()
    out = Path(args.out).resolve()
    pages = out / 'pages'
    assets = out / 'assets'
    
    pages.mkdir(parents=True, exist_ok=True)
    assets.mkdir(parents=True, exist_ok=True)

    items = []
    count = 0

    for root, dirs, files in os.walk(src):
        # Skip generated folder itself
        if 'generated' in root:
            continue
        for f in files:
            if Path(f).suffix.lower() not in SUPPORTED:
                continue
            p = Path(root) / f
            rel = p.relative_to(src)
            print(f"[{count+1}] {rel}")
            count += 1

            # Copy original
            asset = assets / safe_filename(str(rel))
            asset.parent.mkdir(parents=True, exist_ok=True)
            shutil.copy2(p, asset)

            # Extract
            ext = p.suffix.lower()
            if ext == '.docx':
                title, text = extract_docx(p)
            elif ext == '.pptx':
                title, text = extract_pptx(p)
            elif ext == '.pdf':
                title, text = extract_pdf(p)
            elif ext == '.doc':
                title, text = extract_doc(p)
            else:
                title, text = p.stem, ''

            # Write HTML
            page = pages / (safe_filename(title) + f'_{abs(hash(str(rel)))%10000}.html')
            asset_rel = os.path.relpath(asset, page.parent)
            meta = {
                'Type': ext.upper(),
                'Size': f"{p.stat().st_size} bytes",
                'Date': datetime.now().strftime('%Y-%m-%d')
            }
            write_html(page, title, meta, text, asset_rel)
            items.append((title, os.path.relpath(page, out), os.path.relpath(asset, out)))

    # Write index
    idx = f'''<!DOCTYPE html><html><head><meta charset="utf-8"><meta name="viewport" content="width=device-width">
<title>Document Library</title><style>{CSS}</style></head><body>
<h1>📚 Document Library</h1>
<p>Total: {len(items)} documents</p>
<div class="index">
''' + '\n'.join(f'<div class="item"><a href="{html.escape(page)}">{html.escape(title)}</a> <a href="{html.escape(asset)}" style="font-size:0.9rem">📥</a></div>' 
                for title, page, asset in items) + '''
</div></body></html>'''
    (out / 'index.html').write_text(idx, encoding='utf-8')
    print(f"\n✓ Done! {len(items)} documents → {out}")


if __name__ == '__main__':
    main()
