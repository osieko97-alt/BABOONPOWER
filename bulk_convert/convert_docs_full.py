#!/usr/bin/env python3
"""
Enhanced document converter with optional OCR and .doc support.
Usage:
  python convert_docs_full.py --src <src> --out <out> [--ocr] [--convert-doc]

Notes:
- OCR requires Tesseract and poppler to be installed on the system.
- .doc conversion requires LibreOffice/soffice available on PATH.
"""
import argparse
import os
import shutil
import sys
import html
from pathlib import Path
from datetime import datetime

try:
    from docx import Document as DocxDocument
except Exception:
    DocxDocument = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    from PyPDF2 import PdfReader
except Exception:
    PdfReader = None

# OCR deps
try:
    import pytesseract
    from PIL import Image
    from pdf2image import convert_from_path
except Exception:
    pytesseract = None

SUPPORTED = ['.pdf', '.docx', '.pptx', '.doc']

CSS = '''
body{font-family:Segoe UI,Arial,Helvetica,sans-serif;margin:24px}
header{border-bottom:1px solid #ddd;margin-bottom:12px;padding-bottom:8px}
.title{font-size:1.6rem;font-weight:700}
.meta{color:#666;font-size:0.9rem}
.content{white-space:pre-wrap;margin-top:16px}
.item{margin:8px 0;padding:10px;border:1px solid #f0f0f0;border-radius:6px}
'''


def safe_filename(s: str) -> str:
    keep = "-_ .()[]" + "0123456789abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ"
    s2 = ''.join(c if c in keep else '_' for c in s)
    s2 = s2.strip(' _.-')
    if not s2:
        s2 = 'doc'
    return s2


def try_convert_doc_to_docx(doc_path: Path) -> Path:
    # Attempt to convert .doc -> .docx using libreoffice (soffice)
    if doc_path.suffix.lower() != '.doc':
        return doc_path
    out_dir = doc_path.parent / (doc_path.stem + '_converted')
    out_dir.mkdir(parents=True, exist_ok=True)
    soffice = shutil.which('soffice') or shutil.which('libreoffice')
    if not soffice:
        return doc_path
    cmd = f'"{soffice}" --headless --convert-to docx --outdir "{out_dir}" "{doc_path}"'
    print('Running:', cmd)
    res = os.system(cmd)
    if res != 0:
        print('LibreOffice conversion failed (exit', res, ')')
        return doc_path
    converted = out_dir / (doc_path.stem + '.docx')
    if converted.exists():
        return converted
    return doc_path


def extract_text_docx(path: Path) -> (str, str):
    if DocxDocument is None:
        return '', ''
    try:
        doc = DocxDocument(path)
        props = doc.core_properties
        title = props.title or path.stem
        parts = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                parts.append(text)
        return title, '\n\n'.join(parts)
    except Exception as e:
        return path.stem, f"[Error extracting DOCX: {e}]"


def extract_text_pptx(path: Path) -> (str, str):
    if Presentation is None:
        return '', ''
    try:
        prs = Presentation(str(path))
        props = prs.core_properties
        title = props.title or path.stem
        slides = []
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    t = shape.text.strip()
                    if t:
                        slide_text.append(t)
            if slide_text:
                slides.append('\n'.join(slide_text))
        content = '\n\n'.join(f"Slide {i+1}:\n{slides[i]}" for i in range(len(slides)))
        return title, content
    except Exception as e:
        return path.stem, f"[Error extracting PPTX: {e}]"


def extract_text_pdf(path: Path, use_ocr=False) -> (str, str):
    # First try text extraction from PyPDF2
    text_parts = []
    title = path.stem
    if PdfReader is not None:
        try:
            reader = PdfReader(str(path))
            info = reader.metadata
            if info:
                title = info.title or title
            for p in reader.pages:
                try:
                    t = p.extract_text()
                except Exception:
                    t = None
                if t:
                    text_parts.append(t)
        except Exception as e:
            text_parts.append(f"[Error reading PDF: {e}]")

    content = '\n\n'.join(text_parts)
    # If empty and OCR requested and pytesseract available, run OCR
    if (not content.strip()) and use_ocr and pytesseract is not None:
        print('No text found — running OCR on', path)
        try:
            # convert pages to images
            images = convert_from_path(str(path))
            ocr_pages = []
            for i, img in enumerate(images):
                try:
                    txt = pytesseract.image_to_string(img)
                except Exception as e:
                    txt = f'[OCR error on page {i+1}: {e}]'
                ocr_pages.append(txt)
            content = '\n\n'.join(ocr_pages)
        except Exception as e:
            content = f'[PDF OCR failed: {e}]'
    return title, content


def write_html(out_path: Path, title: str, meta: dict, content: str, original_rel: str):
    safe_title = html.escape(title)
    meta_lines = '\n'.join(f"<div><strong>{html.escape(k)}:</strong> {html.escape(str(v))}</div>" for k, v in meta.items() if v)
    body = f"""
<!doctype html>
<html lang="en">
<head>
<meta charset="utf-8" />
<meta name="viewport" content="width=device-width,initial-scale=1" />
<title>{safe_title}</title>
<style>{CSS}</style>
</head>
<body>
<header>
<div class="title">{safe_title}</div>
<div class="meta">{meta_lines}</div>
<div class="meta">Original file: <a href="{html.escape(original_rel)}">{html.escape(os.path.basename(original_rel))}</a></div>
</header>
<main class="content"><pre>{html.escape(content)}</pre></main>
</body>
</html>
"""
    out_path.write_text(body, encoding='utf-8')


def scan_and_convert(src: Path, out: Path, use_ocr=False, convert_doc=False):
    src = src.resolve()
    out = out.resolve()
    assets = out / 'assets'
    pages = out / 'pages'
    out.mkdir(parents=True, exist_ok=True)
    assets.mkdir(parents=True, exist_ok=True)
    pages.mkdir(parents=True, exist_ok=True)

    index_items = []

    for root, dirs, files in os.walk(src):
        for f in files:
            p = Path(root) / f
            ext = p.suffix.lower()
            if ext not in SUPPORTED:
                continue
            rel_src = p.relative_to(src)
            print(f"Processing {rel_src}")

            work_path = p
            if ext == '.doc' and convert_doc:
                converted = try_convert_doc_to_docx(p)
                if converted != p and converted.exists():
                    work_path = converted
                    ext = '.docx'

            # copy original into assets
            target_asset = assets / safe_filename(str(rel_src))
            try:
                shutil.copy2(p, target_asset)
            except Exception:
                target_asset = assets / safe_filename(p.name)
                shutil.copy2(p, target_asset)

            if ext == '.docx':
                title, content = extract_text_docx(work_path)
            elif ext == '.pptx':
                title, content = extract_text_pptx(work_path)
            elif ext == '.pdf':
                title, content = extract_text_pdf(work_path, use_ocr=use_ocr)
            else:
                title, content = work_path.stem, '[Unsupported file type]'

            if not content or not content.strip():
                content = '[No text extracted]'

            meta = {
                'source_path': str(rel_src),
                'converted_at': datetime.utcnow().isoformat() + 'Z',
                'size_bytes': p.stat().st_size,
            }
            page_name = safe_filename(title) + '_' + str(abs(hash(str(rel_src))))[:8] + '.html'
            page_path = pages / page_name
            original_rel = os.path.relpath(target_asset, page_path.parent)
            write_html(page_path, title, meta, content, original_rel)

            index_items.append({
                'title': title,
                'page': os.path.relpath(page_path, out),
                'original': os.path.relpath(target_asset, out),
                'snippet': (content or '')[:300].replace('\n', ' ')
            })

    # write index
    idx_lines = ["<!doctype html>", "<html><head><meta charset=\"utf-8\"><title>Index</title>", f"<style>{CSS}</style>", "</head><body>", "<h1>Documents</h1>", "<div>"]
    for it in index_items:
        idx_lines.append(f"<div class='item'><a href='{html.escape(it['page'])}'><strong>{html.escape(it['title'])}</strong></a><div class='meta'>Original: <a href='{html.escape(it['original'])}'>{html.escape(os.path.basename(it['original']))}</a></div><div>{html.escape(it['snippet'])}...</div></div>")
    idx_lines.append("</div></body></html>")
    (out / 'index.html').write_text('\n'.join(idx_lines), encoding='utf-8')
    print(f"Wrote {len(index_items)} pages to {out}")


def main():
    parser = argparse.ArgumentParser(description='Convert documents to HTML pages (full)')
    parser.add_argument('--src', help='Source folder to scan', default='.')
    parser.add_argument('--out', help='Output folder', default='generated')
    parser.add_argument('--ocr', help='Enable OCR for PDFs with no extracted text', action='store_true')
    parser.add_argument('--convert-doc', help='Attempt to convert .doc files via LibreOffice', action='store_true')
    args = parser.parse_args()

    src = Path(args.src)
    out = Path(args.out)
    if not src.exists():
        print('Source folder does not exist:', src)
        sys.exit(1)

    scan_and_convert(src, out, use_ocr=args.ocr, convert_doc=args.convert_doc)


if __name__ == '__main__':
    main()
